"""Generate the professional Persian Abdit pitch deck.

Run:
    python -m tools.generate_presentation

Outputs AbDit-Presentation.pptx in the repo root.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x07, 0x4D, 0x6E)
PRIMARY = RGBColor(0x0B, 0x6E, 0x9C)
ACCENT = RGBColor(0x18, 0xB6, 0xC9)
LIGHT_BG = RGBColor(0xF3, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x1F, 0x2D, 0x3D)
MUTED = RGBColor(0x6C, 0x7A, 0x89)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x5C, 0xB8, 0x5C)
WARNING = RGBColor(0xF0, 0xAD, 0x4E)
DANGER = RGBColor(0xD9, 0x53, 0x4F)

FONT = "Tahoma"  # Persian-safe fallback present on Windows; replace with Vazirmatn if installed.

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_rtl(paragraph) -> None:
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def style_run(run, *, size=18, bold=False, color=DARK_TEXT, font=FONT) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(text_frame, lines, *, default_size=18, default_color=DARK_TEXT, rtl=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, clear=True):
    """lines: list of dicts/strings.

    Each entry can be a string (uses defaults) or a dict with keys
    text/size/bold/color/level/align.
    """
    if clear:
        text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor

    for i, item in enumerate(lines):
        if isinstance(item, str):
            cfg = {"text": item}
        else:
            cfg = dict(item)
        text = cfg.get("text", "")
        p = text_frame.paragraphs[0] if (i == 0 and clear) else text_frame.add_paragraph()
        p.alignment = cfg.get("align", align)
        p.level = cfg.get("level", 0)
        if rtl:
            set_rtl(p)
        run = p.add_run()
        run.text = text
        style_run(
            run,
            size=cfg.get("size", default_size),
            bold=cfg.get("bold", False),
            color=cfg.get("color", default_color),
            font=cfg.get("font", FONT),
        )
        if cfg.get("space_after"):
            p.space_after = Pt(cfg["space_after"])


def add_textbox(slide, x, y, w, h, lines, **kw):
    box = slide.shapes.add_textbox(x, y, w, h)
    add_text(box.text_frame, lines, **kw)
    return box


def add_rect(slide, x, y, w, h, *, fill=PRIMARY, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    return s


def add_circle(slide, cx, cy, r, *, fill=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def slide_background(slide, color):
    # Fill the slide background by adding a full-bleed rectangle behind everything
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Move to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_chrome(slide, page_num: int, total: int, title: str | None = None):
    """Top accent bar, page number, and tiny brand mark."""
    # Top accent bar (right-aligned thin strip in RTL feel)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Brand wordmark — top right (RTL)
    add_textbox(
        slide,
        Inches(11.5), Inches(0.25),
        Inches(1.7), Inches(0.4),
        [{"text": "💧 آبدیت", "size": 14, "bold": True, "color": NAVY}],
        rtl=True, align=PP_ALIGN.RIGHT, clear=True,
    )

    # Footer page number — bottom left (since RTL puts main content on right)
    add_textbox(
        slide,
        Inches(0.3), Inches(7.05),
        Inches(2.5), Inches(0.4),
        [{"text": f"{page_num} / {total}", "size": 11, "color": MUTED}],
        rtl=False, align=PP_ALIGN.LEFT, clear=True,
    )

    # Footer tagline — bottom right
    add_textbox(
        slide,
        Inches(8.5), Inches(7.05),
        Inches(4.5), Inches(0.4),
        [{"text": "آبدیت · سامانهٔ هوشمند بهینه‌سازی مصرف آب", "size": 11, "color": MUTED}],
        rtl=True, align=PP_ALIGN.RIGHT, clear=True,
    )


def add_slide_title(slide, title: str, subtitle: str | None = None):
    add_textbox(
        slide,
        Inches(0.5), Inches(0.55),
        Inches(12.3), Inches(0.7),
        [{"text": title, "size": 32, "bold": True, "color": NAVY}],
        rtl=True, align=PP_ALIGN.RIGHT,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.5), Inches(1.15),
            Inches(12.3), Inches(0.45),
            [{"text": subtitle, "size": 16, "color": MUTED}],
            rtl=True, align=PP_ALIGN.RIGHT,
        )
    # Title underline accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.7), Inches(1.2), Inches(1.2), Emu(45720))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False


# ---------------------------------------------------------------------------
# Build the deck
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]
    slides_meta: list = []  # (builder, title-or-None)

    def add(builder, title=None):
        slides_meta.append((builder, title))

    # -----------------------------------------------------------------------
    # Slide 1 — Title cover
    # -----------------------------------------------------------------------
    def s1(slide, page, total):
        slide_background(slide, NAVY)
        # Decorative accent shapes
        a1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
        a1.fill.solid(); a1.fill.fore_color.rgb = PRIMARY
        a1.line.fill.background(); a1.shadow.inherit = False
        a2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4.5), Inches(6), Inches(6))
        a2.fill.solid(); a2.fill.fore_color.rgb = ACCENT
        a2.line.fill.background(); a2.shadow.inherit = False

        add_textbox(
            slide,
            Inches(1), Inches(2.2),
            Inches(11.3), Inches(1.3),
            [{"text": "💧 آبدیت", "size": 80, "bold": True, "color": WHITE}],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            Inches(1), Inches(3.6),
            Inches(11.3), Inches(0.8),
            [{"text": "سامانهٔ هوشمند بهینه‌سازی مصرف آب", "size": 30, "color": WHITE}],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            Inches(1), Inches(4.6),
            Inches(11.3), Inches(0.8),
            [{"text": "آب کمتر · آگاهی بیشتر · آیندهٔ پایدارتر", "size": 22, "color": ACCENT}],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            Inches(1), Inches(6.4),
            Inches(11.3), Inches(0.5),
            [{"text": "نسخهٔ نمایشی برای ارائهٔ سرمایه‌گذاران و شرکت‌های آب", "size": 14, "color": WHITE}],
            rtl=True, align=PP_ALIGN.CENTER,
        )

    add(s1)

    # -----------------------------------------------------------------------
    # Slide 2 — Water challenge
    # -----------------------------------------------------------------------
    def s2(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "چالش بحران آب", "چرا اکنون باید اقدام کنیم؟")

        stats = [
            ("۲ میلیارد+", "نفر در جهان با تنش آبی روبه‌رویند", PRIMARY),
            ("۲۰ تا ۴۰٪", "از مصرف خانگی غیرضروری است", ACCENT),
            ("+۸۵٪", "افزایش تعرفه‌ها در دو دههٔ اخیر", DANGER),
        ]
        x0 = Inches(0.6)
        for i, (big, small, color) in enumerate(stats):
            x = x0 + Inches(4.1 * i)
            card = add_rect(slide, x, Inches(2.3), Inches(3.95), Inches(2.6), fill=WHITE,
                            line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            # accent stripe at top
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.3), Inches(3.95), Inches(0.18))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x, Inches(2.7), Inches(3.95), Inches(1.0),
                        [{"text": big, "size": 50, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, x + Inches(0.2), Inches(3.85), Inches(3.55), Inches(1.0),
                        [{"text": small, "size": 16, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.CENTER)

        add_textbox(
            slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.2),
            [
                {"text": "آب در بسیاری از مناطق به یک دارایی استراتژیک تبدیل شده است.",
                 "size": 18, "color": DARK_TEXT},
                {"text": "مدیریت تقاضا، اکنون به اندازهٔ مدیریت عرضه اهمیت دارد.",
                 "size": 16, "color": MUTED, "space_after": 0},
            ],
            rtl=True, align=PP_ALIGN.RIGHT,
        )

    add(s2)

    # -----------------------------------------------------------------------
    # Slide 3 — Traditional model
    # -----------------------------------------------------------------------
    def s3(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "مدل سنتی و محدودیت‌های آن", "صورت‌حساب‌محور، نه راهکارمحور")

        # Flow diagram: 3 boxes with arrows in RTL order (right-to-left)
        labels = ["قرائت کنتور", "صدور صورت‌حساب", "پرداخت توسط مشترک"]
        box_w = Inches(2.6); box_h = Inches(1.0)
        gap = Inches(0.6)
        total_w = box_w * 3 + gap * 2
        start_x = (SLIDE_W - total_w) / 2
        y = Inches(2.4)
        for i, lbl in enumerate(labels):
            x = start_x + (box_w + gap) * i
            r = add_rect(slide, x, y, box_w, box_h, fill=PRIMARY)
            tf = r.text_frame
            tf.margin_left = tf.margin_right = Inches(0.1)
            add_text(tf, [{"text": lbl, "size": 18, "bold": True, "color": WHITE}],
                     rtl=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < 2:
                # Arrow pointing to next (RTL: arrow points left, next is to the right of next-loop iter)
                ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + box_w, y + Inches(0.35), gap, Inches(0.3))
                ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
                ar.line.fill.background(); ar.shadow.inherit = False

        # Limitations
        add_textbox(slide, Inches(0.6), Inches(4.0), Inches(12.1), Inches(0.5),
                    [{"text": "محدودیت‌ها", "size": 22, "bold": True, "color": NAVY}],
                    rtl=True, align=PP_ALIGN.RIGHT)
        bullets = [
            "بدون تحلیل الگوی مصرف",
            "بدون توصیه به مشترک برای کاهش مصرف",
            "هدررفت آب اغلب نامرئی باقی می‌ماند",
            "شرکت‌های آب فاقد ابزار مدیریت تقاضا هستند",
        ]
        add_textbox(
            slide, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.2),
            [{"text": f"•  {b}", "size": 18, "color": DARK_TEXT, "space_after": 8} for b in bullets],
            rtl=True, align=PP_ALIGN.RIGHT,
        )

    add(s3)

    # -----------------------------------------------------------------------
    # Slide 4 — Introducing Abdit
    # -----------------------------------------------------------------------
    def s4(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "معرفی آبدیت", "پلتفرم واسط هوشمند برای کل اکوسیستم آب")

        # Big intro statement
        add_textbox(
            slide, Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.4),
            [
                {"text": "آبدیت چیست؟", "size": 22, "bold": True, "color": NAVY, "space_after": 6},
                {"text":
                    "پلتفرمی که الگوی مصرف آب را تحلیل می‌کند و مشترکان را به متخصصان، "
                    "تأمین‌کنندگان تجهیزات و مؤسسات آموزشی متصل می‌کند تا با راهکارهای واقعی، "
                    "مصرف و هزینه را کاهش دهند.",
                 "size": 17, "color": DARK_TEXT, "space_after": 6},
                {"text":
                    "همان‌گونه که پلتفرم‌های اشتراک سفر رانندگان و مسافران را به هم می‌رسانند، "
                    "آبدیت همهٔ ذی‌نفعان آب را در یک شبکه گرد می‌آورد.",
                 "size": 15, "color": MUTED},
            ],
            rtl=True, align=PP_ALIGN.RIGHT,
        )

        # Four pillars
        pillars = [
            ("🔍", "تحلیل", "تشخیص الگوهای پرمصرف", PRIMARY),
            ("💡", "توصیه", "راهکارهای متناسب با هر خانوار", ACCENT),
            ("🛠", "تجهیز", "بازارگاه محصولات بهینه‌مصرف", SUCCESS),
            ("📚", "یادگیری", "بهبود مداوم با هر پروژه", WARNING),
        ]
        x0 = Inches(0.6); y = Inches(4.0)
        card_w = Inches(2.95); card_h = Inches(2.4)
        gap = Inches(0.2)
        for i, (icon, title, desc, color) in enumerate(pillars):
            x = x0 + (card_w + gap) * i
            add_rect(slide, x, y, card_w, card_h, fill=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x, y + Inches(0.3), card_w, Inches(0.8),
                        [{"text": icon, "size": 36}], rtl=False, align=PP_ALIGN.CENTER)
            add_textbox(slide, x, y + Inches(1.1), card_w, Inches(0.5),
                        [{"text": title, "size": 22, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, x + Inches(0.15), y + Inches(1.6), card_w - Inches(0.3), Inches(0.7),
                        [{"text": desc, "size": 14, "color": MUTED}],
                        rtl=True, align=PP_ALIGN.CENTER)

    add(s4)

    # -----------------------------------------------------------------------
    # Slide 5 — Value proposition by stakeholder
    # -----------------------------------------------------------------------
    def s5(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "ارزش پیشنهادی برای هر ذی‌نفع",
                        "چه چیزی هر طرف را به آبدیت متصل می‌کند")

        rows = [
            ("مشترکان", "کاهش هزینه، آگاهی از الگوی مصرف، دسترسی به فناوری‌های روز", PRIMARY),
            ("شرکت‌های آب", "مدیریت تقاضا، کاهش فشار بر زیرساخت، تعامل بهتر با مشترک", ACCENT),
            ("متخصصان", "بازار کار جدید، ابزار مدیریت پروژه، اعتبار حرفه‌ای", SUCCESS),
            ("تأمین‌کنندگان", "دسترسی به تقاضای هدفمند، بازار جدید، داده‌های عملکرد", WARNING),
            ("مؤسسات آموزشی", "برنامه‌های گواهی‌نامه، بازار شغلی برای فارغ‌التحصیلان", DANGER),
        ]
        y0 = Inches(2.0); h = Inches(0.85); gap = Inches(0.12)
        for i, (who, what, color) in enumerate(rows):
            y = y0 + (h + gap) * i
            # Accent left bar (right side in RTL)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(12.55), y, Inches(0.18), h)
            bar.fill.solid(); bar.fill.fore_color.rgb = color
            bar.line.fill.background(); bar.shadow.inherit = False
            add_rect(slide, Inches(0.6), y, Inches(11.95), h, fill=WHITE)
            add_textbox(slide, Inches(9.0), y + Inches(0.08), Inches(3.5), Inches(0.7),
                        [{"text": who, "size": 20, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, Inches(0.8), y + Inches(0.18), Inches(8.0), Inches(0.6),
                        [{"text": what, "size": 15, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s5)

    # -----------------------------------------------------------------------
    # Slide 6 — 7-step workflow overview
    # -----------------------------------------------------------------------
    def s6(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "روند کار آبدیت — هفت گام", "از داده تا صرفه‌جویی تأییدشده")

        steps = [
            ("۱", "گردآوری\nداده"),
            ("۲", "تحلیل\nمصرف"),
            ("۳", "اعلان به\nمشترک"),
            ("۴", "بازرسی\nمتخصص"),
            ("۵", "پیش‌فاکتور\nو پیشنهاد"),
            ("۶", "نصب\nتجهیزات"),
            ("۷", "بازخورد\nو یادگیری"),
        ]
        # Two rows of cards in RTL order
        card_w = Inches(1.6); card_h = Inches(1.55)
        gap = Inches(0.2)
        total_w = card_w * 7 + gap * 6
        x0 = (SLIDE_W - total_w) / 2
        y = Inches(2.5)
        for i, (num, label) in enumerate(steps):
            x = x0 + (card_w + gap) * i
            add_circle(slide, x + card_w / 2, y + Inches(0.4), Inches(0.4), fill=PRIMARY)
            add_textbox(slide, x, y, card_w, Inches(0.8),
                        [{"text": num, "size": 28, "bold": True, "color": WHITE}],
                        rtl=False, align=PP_ALIGN.CENTER)
            add_rect(slide, x, y + Inches(0.85), card_w, card_h - Inches(0.85), fill=WHITE)
            add_textbox(slide, x + Inches(0.05), y + Inches(0.95), card_w - Inches(0.1), Inches(1.0),
                        [{"text": label, "size": 14, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.CENTER)

        # Loop arrow back from 7 to 1 to show learning loop
        add_textbox(
            slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.5),
            [{"text": "هر پروژهٔ تأییدشده، دانش را به موتور تحلیل بازمی‌گرداند تا توصیه‌های آینده دقیق‌تر شوند.",
              "size": 16, "color": MUTED}],
            rtl=True, align=PP_ALIGN.CENTER,
        )

        # Three call-out blocks under the workflow
        callouts = [
            ("ورودی", "قرائت کنتور · ویژگی‌های خانوار · دادهٔ منطقه‌ای", PRIMARY),
            ("هستهٔ تحلیل", "آمار · یادگیری ماشین · پایگاه دانش", ACCENT),
            ("خروجی", "هشدار به مشترک · بازرسی · توصیهٔ تجهیز · صرفه‌جویی واقعی", SUCCESS),
        ]
        cw = Inches(3.95); ch = Inches(1.1); cy = Inches(5.4); cx0 = Inches(0.6)
        for i, (title, body, color) in enumerate(callouts):
            x = cx0 + (cw + Inches(0.15)) * i
            add_rect(slide, x, cy, cw, ch, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, cy, cw, Inches(0.12))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x + Inches(0.15), cy + Inches(0.18), cw - Inches(0.3), Inches(0.4),
                        [{"text": title, "size": 16, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, x + Inches(0.15), cy + Inches(0.55), cw - Inches(0.3), Inches(0.6),
                        [{"text": body, "size": 13, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s6)

    # -----------------------------------------------------------------------
    # Slide 7 — Steps 1-3 detail (data → analysis → notification)
    # -----------------------------------------------------------------------
    def s7(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "گام‌های ۱ تا ۳ — از داده تا اعلان",
                        "تحلیل خودکار، شناسایی ناهنجاری، اطلاع‌رسانی به مشترک")
        items = [
            ("۱  گردآوری داده", PRIMARY,
             "قرائت‌های کنتور، تاریخچهٔ صورت‌حساب‌ها، موقعیت مکانی و روند مصرف از شرکت آب وارد موتور تحلیل می‌شود."),
            ("۲  تحلیل مصرف", ACCENT,
             "ترکیبی از مدل‌های آماری، یادگیری ماشین و مقایسه با خانوارهای مشابه؛ هر مصرف غیرعادی با درجهٔ شدت کم/متوسط/زیاد علامت‌گذاری می‌شود."),
            ("۳  اعلان به مشترک", SUCCESS,
             "مشترک هشدار همراه با توضیح علمی الگو دریافت می‌کند و می‌تواند با یک کلیک، درخواست بازرسی رایگان دهد. مشارکت کاملاً اختیاری است."),
        ]
        y0 = Inches(2.1); h = Inches(1.5); gap = Inches(0.25)
        for i, (title, color, body) in enumerate(items):
            y = y0 + (h + gap) * i
            add_rect(slide, Inches(0.6), y, Inches(12.1), h, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), y, Inches(0.22), h)
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, Inches(0.8), y + Inches(0.2), Inches(11.6), Inches(0.5),
                        [{"text": title, "size": 22, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, Inches(0.8), y + Inches(0.75), Inches(11.6), Inches(0.7),
                        [{"text": body, "size": 15, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s7)

    # -----------------------------------------------------------------------
    # Slide 8 — Steps 4-7 detail (inspection → install → learn)
    # -----------------------------------------------------------------------
    def s8(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "گام‌های ۴ تا ۷ — از بازرسی تا یادگیری",
                        "اجرای واقعی راهکار و بازخورد به سامانه")
        items = [
            ("۴  بازرسی متخصص", PRIMARY,
             "متخصص دارای گواهی به محل می‌رود؛ لوله‌کشی، آبیاری، مصرف داخلی و نشت‌ها را بررسی می‌کند."),
            ("۵  پیش‌فاکتور و پیشنهاد", ACCENT,
             "گزارش فنی همراه با محصولات پیشنهادی از بازارگاه و برآورد صرفه‌جویی برای مشترک ارسال می‌شود."),
            ("۶  نصب تجهیزات", SUCCESS,
             "در صورت پذیرش، تأمین‌کنندگان تجهیزات را تأمین کرده و متخصصان دارای گواهی نصب می‌کنند."),
            ("۷  بازخورد و یادگیری", WARNING,
             "صرفه‌جویی واقعی پس از نصب اندازه‌گیری و به پایگاه دانش بازگردانده می‌شود تا توصیه‌های آینده دقیق‌تر شوند."),
        ]
        y0 = Inches(2.1); h = Inches(1.1); gap = Inches(0.18)
        for i, (title, color, body) in enumerate(items):
            y = y0 + (h + gap) * i
            add_rect(slide, Inches(0.6), y, Inches(12.1), h, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), y, Inches(0.22), h)
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, Inches(0.8), y + Inches(0.12), Inches(11.6), Inches(0.45),
                        [{"text": title, "size": 18, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, Inches(0.8), y + Inches(0.55), Inches(11.6), Inches(0.5),
                        [{"text": body, "size": 14, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s8)

    # -----------------------------------------------------------------------
    # Slide 9 — Stakeholders (radial)
    # -----------------------------------------------------------------------
    def s9(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "ذی‌نفعان آبدیت", "پنج طرف، یک شبکه")

        cx, cy = Inches(6.66), Inches(4.5)
        # Center
        add_circle(slide, cx, cy, Inches(0.95), fill=NAVY)
        add_textbox(slide, cx - Inches(0.95), cy - Inches(0.35), Inches(1.9), Inches(0.7),
                    [{"text": "آبدیت", "size": 22, "bold": True, "color": WHITE}],
                    rtl=True, align=PP_ALIGN.CENTER)

        roles = [
            ("مشترکان", PRIMARY, Inches(2.4), Inches(2.6)),
            ("شرکت‌های آب", ACCENT, Inches(10.9), Inches(2.6)),
            ("متخصصان", SUCCESS, Inches(1.5), Inches(5.4)),
            ("تأمین‌کنندگان", WARNING, Inches(11.7), Inches(5.4)),
            ("مؤسسات آموزشی", DANGER, Inches(6.66) - Inches(1.0), Inches(6.4)),
        ]
        for label, color, x, y in roles:
            r = Inches(0.9)
            add_circle(slide, x, y, r, fill=color)
            add_textbox(slide, x - r, y - Inches(0.3), r * 2, Inches(0.6),
                        [{"text": label, "size": 14, "bold": True, "color": WHITE}],
                        rtl=True, align=PP_ALIGN.CENTER)
            # connector line
            ln = slide.shapes.add_connector(1, x, y, cx, cy)
            ln.line.color.rgb = MUTED
            ln.line.width = Pt(1.2)

    add(s9)

    # -----------------------------------------------------------------------
    # Slide 10 — System components
    # -----------------------------------------------------------------------
    def s10(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "اجزای فنی سامانه",
                        "موتور تحلیل، پایگاه دانش، پورتال‌های نقش‌محور و بازارگاه")
        cards = [
            ("موتور تحلیل داده", "تشخیص الگو · ناهنجاری · مقایسهٔ همتا"),
            ("پایگاه دانش", "راهکارها · مطالعات موردی · بهترین تجربه‌های منطقه‌ای"),
            ("پورتال مشترک", "نمای مصرف · اعلان · درخواست بازرسی · پاسخ به پیش‌فاکتور"),
            ("پورتال متخصص", "دریافت کار · ثبت گزارش · ساخت پیشنهاد · پیگیری نصب"),
            ("بازارگاه تأمین‌کنندگان", "ثبت محصول · قیمت‌گذاری · ادغام با خدمات نصب"),
            ("داشبورد شرکت آب", "روند مصرف · ناهنجاری‌ها · مدیریت تقاضا"),
        ]
        cw = Inches(4.0); ch = Inches(1.7); gx = Inches(0.13); gy = Inches(0.18)
        x0 = (SLIDE_W - (cw * 3 + gx * 2)) / 2
        y0 = Inches(2.1)
        for i, (title, body) in enumerate(cards):
            row, col = divmod(i, 3)
            x = x0 + (cw + gx) * col
            y = y0 + (ch + gy) * row
            add_rect(slide, x, y, cw, ch, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.12))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = PRIMARY
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x + Inches(0.2), y + Inches(0.18), cw - Inches(0.4), Inches(0.5),
                        [{"text": title, "size": 18, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(1.0),
                        [{"text": body, "size": 14, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s10)

    # -----------------------------------------------------------------------
    # Slide 11 — Marketplace
    # -----------------------------------------------------------------------
    def s11(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "بازارگاه تجهیزات بهینه‌مصرف",
                        "شش دستهٔ کلیدی، تأمین‌کنندگان متعدد، کیفیت تأییدشده")

        cats = [
            ("🌱", "آبیاری هوشمند", "کنترلر مبتنی بر شرایط جوی · حسگر رطوبت خاک"),
            ("🚿", "شیرآلات و دوش کم‌مصرف", "هوادهنده · دوش‌های ۴-۶ لیتری در دقیقه"),
            ("🛡", "نشت‌یاب", "تشخیص صوتی · قطع خودکار جریان"),
            ("♻️", "بازچرخانی آب خاکستری", "استفادهٔ مجدد در فلاش و آبیاری"),
            ("📡", "کنتور هوشمند", "گزارش لحظه‌ای · هشدار نشت"),
            ("🧪", "کنترل کیفیت آب", "حسگرهای کیفیت · پایش پیوسته"),
        ]
        cw = Inches(4.0); ch = Inches(1.65); gx = Inches(0.13); gy = Inches(0.2)
        x0 = (SLIDE_W - (cw * 3 + gx * 2)) / 2
        y0 = Inches(2.1)
        for i, (icon, name, desc) in enumerate(cats):
            row, col = divmod(i, 3)
            x = x0 + (cw + gx) * col
            y = y0 + (ch + gy) * row
            add_rect(slide, x, y, cw, ch, fill=WHITE)
            add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.8),
                        [{"text": icon, "size": 36}], rtl=False, align=PP_ALIGN.LEFT)
            add_textbox(slide, x + Inches(1.1), y + Inches(0.18), cw - Inches(1.3), Inches(0.5),
                        [{"text": name, "size": 18, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, x + Inches(0.2), y + Inches(0.95), cw - Inches(0.4), Inches(0.7),
                        [{"text": desc, "size": 13, "color": MUTED}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s11)

    # -----------------------------------------------------------------------
    # Slide 12 — Learning loop
    # -----------------------------------------------------------------------
    def s12(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "حلقهٔ یادگیری مداوم",
                        "هر پروژه، توصیه‌های آینده را دقیق‌تر می‌کند")

        nodes = [
            ("داده", "قرائت کنتور · ویژگی‌های خانوار", PRIMARY, Inches(2.0), Inches(2.6)),
            ("تحلیل", "آمار · یادگیری ماشین · پایگاه دانش", ACCENT, Inches(7.0), Inches(2.6)),
            ("اقدام", "بازرسی · توصیه · نصب", SUCCESS, Inches(11.0), Inches(4.5)),
            ("بازخورد", "صرفه‌جویی واقعی پس از نصب", WARNING, Inches(7.0), Inches(6.0)),
            ("به‌روزرسانی دانش", "افزایش دقت توصیه‌ها", DANGER, Inches(2.0), Inches(4.5)),
        ]
        for label, body, color, x, y in nodes:
            r = Inches(1.05)
            add_circle(slide, x, y, r, fill=color)
            add_textbox(slide, x - r, y - Inches(0.4), r * 2, Inches(0.5),
                        [{"text": label, "size": 16, "bold": True, "color": WHITE}],
                        rtl=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, x - r, y, r * 2, Inches(0.7),
                        [{"text": body, "size": 11, "color": WHITE}],
                        rtl=True, align=PP_ALIGN.CENTER)

        # Connecting lines (simple polyline approximation via individual connectors)
        pairs = [
            (Inches(2.0), Inches(2.6), Inches(7.0), Inches(2.6)),
            (Inches(7.0), Inches(2.6), Inches(11.0), Inches(4.5)),
            (Inches(11.0), Inches(4.5), Inches(7.0), Inches(6.0)),
            (Inches(7.0), Inches(6.0), Inches(2.0), Inches(4.5)),
            (Inches(2.0), Inches(4.5), Inches(2.0), Inches(2.6)),
        ]
        for x1, y1, x2, y2 in pairs:
            ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
            ln.line.color.rgb = MUTED
            ln.line.width = Pt(2)

    add(s12)

    # -----------------------------------------------------------------------
    # Slide 13 — Sample dashboards (text mock)
    # -----------------------------------------------------------------------
    def s13(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "داشبوردهای نمونه — نسخهٔ نمایشی زنده",
                        "آنچه هر نقش هنگام ورود می‌بیند")

        # Two big mock boxes
        x_left = Inches(0.6); x_right = Inches(7.0); w = Inches(5.7); h = Inches(4.7)
        # Subscriber mock
        add_rect(slide, x_right, Inches(2.0), w, h, fill=WHITE)
        add_textbox(slide, x_right + Inches(0.2), Inches(2.1), w - Inches(0.4), Inches(0.5),
                    [{"text": "داشبورد مشترک", "size": 18, "bold": True, "color": PRIMARY}],
                    rtl=True, align=PP_ALIGN.RIGHT)
        sub_kpis = [
            ("آخرین صورت‌حساب", "۸۹٫۲۰"),
            ("روند نسبت به ماه قبل", "+۳۲٪"),
            ("مصرف سالانه", "۲۸۸ مترمکعب"),
            ("هشدار خوانده‌نشده", "۱"),
        ]
        for i, (k, v) in enumerate(sub_kpis):
            row, col = divmod(i, 2)
            cx = x_right + Inches(0.3) + Inches(2.6) * col
            cy = Inches(2.7) + Inches(0.95) * row
            add_rect(slide, cx, cy, Inches(2.4), Inches(0.85), fill=LIGHT_BG, line=None)
            add_textbox(slide, cx + Inches(0.1), cy + Inches(0.05), Inches(2.2), Inches(0.35),
                        [{"text": k, "size": 11, "color": MUTED}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, cx + Inches(0.1), cy + Inches(0.4), Inches(2.2), Inches(0.5),
                        [{"text": v, "size": 18, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.RIGHT)
        add_textbox(slide, x_right + Inches(0.2), Inches(4.7), w - Inches(0.4), Inches(1.8),
                    [
                        {"text": "نمودار مصرف ۱۲ ماهه + لیست بازرسی‌ها + اعلان‌های فعال",
                         "size": 13, "color": MUTED},
                        {"text": "اقدام یک‌کلیکی: «درخواست بازرسی»",
                         "size": 13, "color": SUCCESS},
                    ],
                    rtl=True, align=PP_ALIGN.RIGHT)

        # Operator mock
        add_rect(slide, x_left, Inches(2.0), w, h, fill=WHITE)
        add_textbox(slide, x_left + Inches(0.2), Inches(2.1), w - Inches(0.4), Inches(0.5),
                    [{"text": "نمای کلی اپراتور", "size": 18, "bold": True, "color": ACCENT}],
                    rtl=True, align=PP_ALIGN.RIGHT)
        op_kpis = [
            ("مشترکان فعال", "۱٬۲۸۰"),
            ("متخصصان", "۲۴"),
            ("ناهنجاری‌های باز", "۸۷"),
            ("میانگین صرفه‌جویی تأییدشده", "۲۳٫۵٪"),
        ]
        for i, (k, v) in enumerate(op_kpis):
            row, col = divmod(i, 2)
            cx = x_left + Inches(0.3) + Inches(2.6) * col
            cy = Inches(2.7) + Inches(0.95) * row
            add_rect(slide, cx, cy, Inches(2.4), Inches(0.85), fill=LIGHT_BG, line=None)
            add_textbox(slide, cx + Inches(0.1), cy + Inches(0.05), Inches(2.2), Inches(0.35),
                        [{"text": k, "size": 11, "color": MUTED}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, cx + Inches(0.1), cy + Inches(0.4), Inches(2.2), Inches(0.5),
                        [{"text": v, "size": 18, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.RIGHT)
        add_textbox(slide, x_left + Inches(0.2), Inches(4.7), w - Inches(0.4), Inches(1.8),
                    [
                        {"text": "روند بازرسی‌ها · صرفه‌جویی تأییدشده · پایگاه دانش زنده",
                         "size": 13, "color": MUTED},
                        {"text": "اجرای دوبارهٔ تحلیل با یک کلیک",
                         "size": 13, "color": SUCCESS},
                    ],
                    rtl=True, align=PP_ALIGN.RIGHT)

    add(s13)

    # -----------------------------------------------------------------------
    # Slide 14 — Expected results
    # -----------------------------------------------------------------------
    def s14(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "نتایج مورد انتظار",
                        "ارقام بر اساس مطالعات موردی و پایگاه دانش آبدیت")

        # Big highlight numbers
        kpis = [
            ("۱۵–۳۰٪", "کاهش معمول مصرف خانوارهای پُرمصرف", PRIMARY),
            ("۲۲٪", "میانگین صرفه‌جویی پس از اصلاح نشت", ACCENT),
            ("< ۶ ماه", "بازگشت سرمایهٔ هوادهنده‌های شیرآلات", SUCCESS),
            ("۲۸٪", "کاهش مصرف آب آشامیدنی با بازچرخانی", WARNING),
        ]
        cw = Inches(2.95); ch = Inches(2.4); gx = Inches(0.18)
        x0 = (SLIDE_W - (cw * 4 + gx * 3)) / 2
        y = Inches(2.5)
        for i, (big, label, color) in enumerate(kpis):
            x = x0 + (cw + gx) * i
            add_rect(slide, x, y, cw, ch, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.18))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x, y + Inches(0.45), cw, Inches(1.0),
                        [{"text": big, "size": 50, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, x + Inches(0.15), y + Inches(1.55), cw - Inches(0.3), Inches(0.8),
                        [{"text": label, "size": 14, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.CENTER)

        add_textbox(
            slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.0),
            [{"text":
                "هر مدخل پایگاه دانش با تعداد پروژه‌های انجام‌شده ارزش‌گذاری می‌شود؛ ارقام با گذشت زمان دقیق‌تر می‌شوند.",
              "size": 14, "color": MUTED}],
            rtl=True, align=PP_ALIGN.CENTER,
        )

    add(s14)

    # -----------------------------------------------------------------------
    # Slide 15 — Business model
    # -----------------------------------------------------------------------
    def s15(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "مدل کسب‌وکار",
                        "چندین جریان درآمد، هم‌راستا با ارزش ایجادشده")

        streams = [
            ("کمیسیون فروش تجهیزات", "درصدی از فروش هر کالای موفق در بازارگاه", PRIMARY),
            ("هزینهٔ بازرسی", "حق‌الزحمهٔ هر بازرسی موفق توسط متخصص", ACCENT),
            ("اشتراک تأمین‌کنندگان", "حضور در بازارگاه · ابزار مدیریت محصول", SUCCESS),
            ("خدمات تحلیل برای شرکت‌های آب", "گزارش مدیریت تقاضا · مشاورهٔ تخصصی", WARNING),
            ("همکاری با مراکز آموزشی", "صدور گواهی متخصصان · تأمین کاندیداها", DANGER),
        ]
        y0 = Inches(2.1); h = Inches(0.85); gap = Inches(0.13)
        for i, (title, body, color) in enumerate(streams):
            y = y0 + (h + gap) * i
            add_rect(slide, Inches(0.6), y, Inches(12.1), h, fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), y, Inches(0.22), h)
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, Inches(8.5), y + Inches(0.15), Inches(4.0), Inches(0.5),
                        [{"text": title, "size": 18, "bold": True, "color": color}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, Inches(0.8), y + Inches(0.18), Inches(7.6), Inches(0.5),
                        [{"text": body, "size": 14, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s15)

    # -----------------------------------------------------------------------
    # Slide 16 — Roadmap
    # -----------------------------------------------------------------------
    def s16(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "نقشهٔ راه آینده", "از نمونهٔ نمایشی تا سامانهٔ کشوری")

        phases = [
            ("اکنون", "نمونهٔ نمایشی کامل با تمام نقش‌ها و گردش‌کار هفت‌مرحله‌ای", PRIMARY),
            ("سه‌ماهه ۱", "ادغام با کنتورهای هوشمند · هشدار آنی نشت", ACCENT),
            ("سه‌ماهه ۲", "مدل پیش‌بینی تقاضا با هوش مصنوعی · نسخهٔ موبایل متخصص", SUCCESS),
            ("سه‌ماهه ۳", "ادغام با خانهٔ هوشمند · پرداخت آنلاین پیشنهادها", WARNING),
            ("سه‌ماهه ۴", "ابزار برنامه‌ریزی منطقه‌ای برای دولت‌های محلی", DANGER),
        ]
        # Horizontal timeline
        x0 = Inches(0.8); y = Inches(2.6)
        cw = Inches(2.4); gap = Inches(0.05)
        for i, (when, what, color) in enumerate(phases):
            x = x0 + (cw + gap) * i
            add_rect(slide, x, y, cw, Inches(2.6), fill=WHITE)
            stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, Inches(0.4))
            stripe.fill.solid(); stripe.fill.fore_color.rgb = color
            stripe.line.fill.background(); stripe.shadow.inherit = False
            add_textbox(slide, x, y + Inches(0.05), cw, Inches(0.4),
                        [{"text": when, "size": 14, "bold": True, "color": WHITE}],
                        rtl=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, x + Inches(0.15), y + Inches(0.6), cw - Inches(0.3), Inches(2.0),
                        [{"text": what, "size": 13, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s16)

    # -----------------------------------------------------------------------
    # Slide 17 — Why Abdit wins
    # -----------------------------------------------------------------------
    def s17(slide, page, total):
        slide_background(slide, LIGHT_BG)
        add_chrome(slide, page, total)
        add_slide_title(slide, "چرا آبدیت برنده می‌شود",
                        "ترکیبی از داده، شبکه، و یادگیری مداوم")

        items = [
            ("از داده تا اقدام", "توصیهٔ مبتنی بر داده، نه پوستر و توصیه‌های کلی"),
            ("شبکهٔ کامل ذی‌نفعان", "همه طرف‌ها در یک پلتفرم — اثر شبکه‌ای واقعی"),
            ("یادگیری مداوم", "هر پروژهٔ جدید، توصیه‌های آینده را بهبود می‌دهد"),
            ("مقیاس‌پذیر و چندشهری", "مدل آماده برای ورود به شهرها و مناطق متعدد"),
            ("شفافیت و اختیار مشترک", "همه‌چیز اختیاری · ارزش روشن قبل از پرداخت"),
        ]
        y0 = Inches(2.2); h = Inches(0.78); gap = Inches(0.12)
        for i, (title, body) in enumerate(items):
            y = y0 + (h + gap) * i
            add_rect(slide, Inches(0.6), y, Inches(12.1), h, fill=WHITE)
            num_circle = add_circle(slide, Inches(12.0), y + Inches(0.39), Inches(0.32), fill=PRIMARY)
            add_textbox(slide, Inches(11.7), y + Inches(0.07), Inches(0.7), Inches(0.6),
                        [{"text": str(i + 1), "size": 18, "bold": True, "color": WHITE}],
                        rtl=False, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(7.5), y + Inches(0.13), Inches(3.8), Inches(0.55),
                        [{"text": title, "size": 18, "bold": True, "color": NAVY}],
                        rtl=True, align=PP_ALIGN.RIGHT)
            add_textbox(slide, Inches(0.8), y + Inches(0.17), Inches(6.5), Inches(0.55),
                        [{"text": body, "size": 14, "color": DARK_TEXT}],
                        rtl=True, align=PP_ALIGN.RIGHT)

    add(s17)

    # -----------------------------------------------------------------------
    # Slide 18 — Closing / call to action
    # -----------------------------------------------------------------------
    def s18(slide, page, total):
        slide_background(slide, NAVY)
        add_chrome(slide, page, total)
        add_textbox(
            slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.4),
            [{"text": "بیایید آب را هوشمند مصرف کنیم", "size": 50, "bold": True, "color": WHITE}],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.8),
            [{"text": "آبدیت آماده است تا با شرکت‌های آب، شهرداری‌ها و سرمایه‌گذاران همکار شود.",
              "size": 22, "color": ACCENT}],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, Inches(1), Inches(4.5), Inches(11.3), Inches(1.0),
            [
                {"text": "نسخهٔ نمایشی زنده: محلی، آمادهٔ اجرا (./run.sh یا run.bat)",
                 "size": 16, "color": WHITE},
                {"text": "همهٔ نقش‌ها · کل گردش‌کار هفت‌مرحله‌ای · داده‌های نمایشی پارسی",
                 "size": 14, "color": WHITE},
            ],
            rtl=True, align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
            [{"text": "ممنون از توجه شما", "size": 22, "bold": True, "color": WHITE}],
            rtl=True, align=PP_ALIGN.CENTER,
        )

    add(s18)

    # -----------------------------------------------------------------------
    # Render every slide
    # -----------------------------------------------------------------------
    total = len(slides_meta)
    for i, (builder, _) in enumerate(slides_meta, start=1):
        slide = prs.slides.add_slide(blank)
        builder(slide, i, total)

    out = Path(__file__).resolve().parent.parent / "AbDit-Presentation.pptx"
    prs.save(out)
    print(f"Wrote {out} ({total} slides)")
    return out


if __name__ == "__main__":
    build()
